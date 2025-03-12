import io
from typing import Dict, Any, List

from langflow.custom import Component
from langflow.inputs.inputs import DataInput
from langflow.io import Output
from langflow.schema import Data, Message


class GetFilesTextContentComponent(Component):
    """Component for extracting text content from different file types."""

    display_name: str = "Get Files Text Content"
    description: str = "Extract text content from files (PDF, DOCX, XLSX, etc)"
    name: str = "GetFilesTextContent"
    icon = "file-text"

    inputs = [
        DataInput(
            name="files",
            display_name="Files",
            info="Files data from download_files component",
            required=True,
        ),
    ]

    outputs = [
        Output(
            name="documents",
            display_name="Documents",
            method="parse_file",
        ),
    ]

    def _parse_pdf(self, content: bytes) -> str:
        """Extract text from PDF file."""
        try:
            from pypdf import PdfReader
            pdf_file = io.BytesIO(content)
            pdf_reader = PdfReader(pdf_file)
            text = ''
            for page in pdf_reader.pages:
                text += page.extract_text() + '\n'
            return text
        except ImportError:
            return "[Error: PDF parsing library not installed]"

    def _clean_text(self, text: str) -> str:
        """Clean up text by removing extra whitespace and normalizing spaces."""
        # Normalize whitespace first
        text = ' '.join(text.split())

        # Fix common word breaks and formatting issues
        text = text.replace(' .', '.')
        text = text.replace(' ,', ',')
        text = text.replace(' :', ':')
        text = text.replace(' ;', ';')
        text = text.replace(' )', ')')
        text = text.replace('( ', '(')

        # Fix broken words (e.g., 'O NLYOFFICE' -> 'ONLYOFFICE')
        common_breaks = [
            ('O NLYOFFICE', 'ONLYOFFICE'),
            ('HTML 5', 'HTML5'),
            ('M S', 'MS'),
            ('Google Docs', 'GoogleDocs'),
            ('Office 365', 'Office365')
        ]
        for broken, fixed in common_breaks:
            text = text.replace(broken, fixed)

        # Remove redundant spaces around special characters
        special_chars = ['-', '/', '\\', '&']
        for char in special_chars:
            text = text.replace(f' {char} ', char)

        return text.strip()

    def _extract_shape_text(self, shape) -> str:
        """Extract text from a shape or inline object."""
        text = ''
        try:
            # Try to get text from shape properties
            if hasattr(shape, 'text'):
                text = shape.text
            elif hasattr(shape, 'text_frame'):
                text = shape.text_frame.text
            # Add more shape text extraction methods as needed
        except Exception:
            pass
        return self._clean_text(text)

    def _parse_run_elements(self, paragraph) -> str:
        """Extract text from paragraph run elements including shapes."""
        text_parts = []
        current_text = ''

        for run in paragraph.runs:
            # Get regular text
            if run.text.strip():
                # If there's a space before and after, preserve it
                if current_text and not current_text.endswith(' ') and not run.text.startswith(' '):
                    current_text += ' '
                current_text += run.text

            # Get text from shapes in the run
            if hasattr(run, '_element'):
                for elem in run._element:
                    if hasattr(elem, 'graphic'):
                        shape_text = self._extract_shape_text(elem)
                        if shape_text:
                            if current_text:
                                text_parts.append(
                                    self._clean_text(current_text))
                                current_text = ''
                            text_parts.append(f"[Shape] {shape_text}")

        if current_text:
            text_parts.append(self._clean_text(current_text))

        return ' '.join(text_parts)

    def _parse_docx(self, content: bytes) -> str:
        """Extract text from DOCX file including all content types."""
        # Try to parse as regular DOCX first
        try:
            from docx import Document
            docx_file = io.BytesIO(content)
            doc = Document(docx_file)
            return self._parse_docx_document(doc)
        except ImportError:
            return "[Error: DOCX parsing library not installed]"
        except Exception:
            # If regular parsing fails, try raw XML parsing
            return self._parse_docx_raw(content)

    def _parse_docx_raw(self, content: bytes) -> str:
        """Parse text directly from DOCX XML structure."""
        import zipfile
        import xml.etree.ElementTree as ET

        text_parts = []

        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zip_ref:
                # Parse main document.xml
                if 'word/document.xml' in zip_ref.namelist():
                    with zip_ref.open('word/document.xml') as doc_xml:
                        tree = ET.parse(doc_xml)
                        root = tree.getroot()

                        # Define namespace
                        ns = {
                            'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

                        # Extract text from paragraphs
                        for para in root.findall('.//w:p', ns):
                            para_text = []
                            for text_elem in para.findall('.//w:t', ns):
                                if text_elem.text:
                                    para_text.append(text_elem.text)
                            if para_text:
                                text_parts.append(' '.join(para_text))

                        # Extract text from text boxes and shapes
                        for shape in root.findall('.//w:drawing//wp:docPr', {'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing'}):
                            if 'title' in shape.attrib:
                                text_parts.append(
                                    f"[Shape: {shape.attrib['title']}]")
                            if 'descr' in shape.attrib:
                                text_parts.append(
                                    f"[Shape Description: {shape.attrib['descr']}]")

        except Exception as e:
            text_parts.append(f"[Error parsing DOCX: {str(e)}]")

        return '\n'.join(text_parts)

    def _parse_docx_document(self, doc) -> str:
        """Extract text from DOCX using python-docx library."""
        text_parts = []

        # Get text from paragraphs including shapes
        for paragraph in doc.paragraphs:
            text = self._parse_run_elements(paragraph)
            if text:
                text_parts.append(text)

        # Get text from tables
        for table in doc.tables:
            for row in table.rows:
                row_texts = []
                for cell in row.cells:
                    cell_text = ' '.join(
                        self._parse_run_elements(p) for p in cell.paragraphs
                        if self._parse_run_elements(p))
                    if cell_text:
                        row_texts.append(cell_text)
                if row_texts:
                    text_parts.append('\t'.join(row_texts))

        # Get text from sections (headers/footers)
        for section in doc.sections:
            # Header
            if section.header:
                header_texts = []
                for paragraph in section.header.paragraphs:
                    text = self._parse_run_elements(paragraph)
                    if text:
                        header_texts.append(text)
                if header_texts:
                    text_parts.append(f"[Header] {' '.join(header_texts)}")

            # Footer
            if section.footer:
                footer_texts = []
                for paragraph in section.footer.paragraphs:
                    text = self._parse_run_elements(paragraph)
                    if text:
                        footer_texts.append(text)
                if footer_texts:
                    text_parts.append(f"[Footer] {' '.join(footer_texts)}")

        # Clean up and join all parts with proper spacing
        cleaned_parts = []
        for part in text_parts:
            if part.strip():
                # Handle special section markers
                if part.startswith('[') and ']' in part:
                    cleaned_parts.append(part)  # Keep section markers as is
                else:
                    # Clean and add only if not empty after cleaning
                    cleaned = self._clean_text(part)
                    if cleaned:
                        cleaned_parts.append(cleaned)

        return '\n'.join(cleaned_parts)

    def _parse_xlsx(self, content: bytes) -> str:
        """Extract text from XLSX file."""
        try:
            import openpyxl
            xlsx_file = io.BytesIO(content)
            workbook = openpyxl.load_workbook(xlsx_file, data_only=True)
            text = []

            for sheet in workbook.worksheets:
                text.append(f"Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    # Filter out None values and convert all to strings
                    row_text = [str(cell) for cell in row if cell is not None]
                    if row_text:  # Only add non-empty rows
                        text.append('\t'.join(row_text))
                text.append('')  # Add blank line between sheets

            return '\n'.join(text)
        except ImportError:
            return "[Error: XLSX parsing library not installed]"
        except Exception as e:
            return f"[Error parsing XLSX: {str(e)}]"

    def _parse_text(self, content: bytes) -> str:
        """Extract text from text-based files."""
        try:
            return content.decode('utf-8')
        except UnicodeDecodeError:
            # Try different encodings if UTF-8 fails
            encodings = ['latin-1', 'cp1252', 'ascii']
            for encoding in encodings:
                try:
                    return content.decode(encoding)
                except UnicodeDecodeError:
                    continue
            raise ValueError(
                "Could not decode text content with any supported encoding")

    def _parse_pptx(self, content: bytes) -> str:
        """Extract text from PPTX file including all content types."""
        try:
            # Try regular parsing first
            from pptx import Presentation
            pptx_file = io.BytesIO(content)
            prs = Presentation(pptx_file)
            return self._parse_pptx_presentation(prs)
        except ImportError:
            return "[Error: PPTX parsing library not installed]"
        except Exception:
            # Fall back to raw XML parsing
            return self._parse_pptx_raw(content)

    def _parse_pptx_presentation(self, prs) -> str:
        """Parse PPTX using python-pptx library."""
        text_parts = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"[Slide {i}]"]

            # Get text from shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    # Handle tables specially
                    if shape.has_table:
                        rows = []
                        for row in shape.table.rows:
                            row_texts = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_texts.append(cell.text.strip())
                            if row_texts:
                                rows.append('\t'.join(row_texts))
                        if rows:
                            slide_parts.append('[Table]\n' + '\n'.join(rows))
                    else:
                        # Regular shape text
                        slide_parts.append(shape.text.strip())

                # Get text from grouped shapes
                if hasattr(shape, 'shapes'):
                    for subshape in shape.shapes:
                        if hasattr(subshape, 'text') and subshape.text.strip():
                            slide_parts.append(subshape.text.strip())

            if len(slide_parts) > 1:  # If we have more than just the slide header
                text_parts.extend(slide_parts)

        return '\n'.join(text_parts)

    def _parse_pptx_raw(self, content: bytes) -> str:
        """Parse text directly from PPTX XML structure."""
        import zipfile
        import xml.etree.ElementTree as ET

        text_parts = []

        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zip_ref:
                # Get list of all slide files
                slide_files = [f for f in zip_ref.namelist(
                ) if f.startswith('ppt/slides/slide')]
                slide_files.sort()

                for i, slide_file in enumerate(slide_files, 1):
                    slide_parts = [f"[Slide {i}]"]

                    with zip_ref.open(slide_file) as slide_xml:
                        tree = ET.parse(slide_xml)
                        root = tree.getroot()

                        # Define namespaces
                        ns = {
                            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
                        }

                        # Extract text from text runs
                        for text_elem in root.findall('.//a:t', ns):
                            if text_elem.text and text_elem.text.strip():
                                slide_parts.append(text_elem.text.strip())

                        # Extract text from notes (if any)
                        notes_file = f'ppt/notesSlides/notesSlide{i}.xml'
                        if notes_file in zip_ref.namelist():
                            with zip_ref.open(notes_file) as notes_xml:
                                notes_tree = ET.parse(notes_xml)
                                notes_root = notes_tree.getroot()
                                notes_texts = []
                                for text_elem in notes_root.findall('.//a:t', ns):
                                    if text_elem.text and text_elem.text.strip():
                                        notes_texts.append(
                                            text_elem.text.strip())
                                if notes_texts:
                                    slide_parts.append(
                                        '[Notes]\n' + '\n'.join(notes_texts))

                    if len(slide_parts) > 1:  # If we have more than just the slide header
                        text_parts.extend(slide_parts)

        except Exception as e:
            text_parts.append(f"[Error parsing PPTX: {str(e)}]")

        return '\n'.join(text_parts)

    def parse_file(self) -> Data:
        """Parse file content based on its type and extract text."""
        try:
            file_data = self.files
            if not isinstance(file_data, Data) or not file_data.data:
                return Data(data={"files": []})

            files: List[Dict[str, Any]] = file_data.data.get('files', [])
            if not files:
                return Data(data={"files": []})

            data = []

            for file in files:
                try:
                    content = file.get('content')
                    if not content:
                        continue

                    extension = file.get('extension', '').lower()

                    # Extract text based on file extension
                    if extension == '.pdf':
                        text = self._parse_pdf(content)
                    elif extension in ['.docx', '.doc']:
                        text = self._parse_docx(content)
                    elif extension in ['.xlsx', '.xls']:
                        text = self._parse_xlsx(content)
                    elif extension in ['.pptx', '.ppt']:
                        text = self._parse_pptx(content)
                    elif extension in ['.txt', '.md', '.py', '.json', '.yaml', '.yml']:
                        text = self._parse_text(content)
                    else:
                        # For unsupported formats, try to parse as text
                        try:
                            text = self._parse_text(content)
                        except ValueError:
                            text = f"[Could not extract text from file with extension {extension}]"

                    item = {
                        'text': text,
                        'extension': extension,
                        'title': file.get('title', ''),
                        'id': file.get('id', ''),
                        'version': file.get('version', ''),
                        'size': len(content)
                    }

                    data.append(item)
                except Exception as e:
                    print(
                        f"Error processing file {file.get('id', 'unknown')}: {str(e)}")
                    continue

            return Data(data={"files": data})

        except Exception as e:
            raise ValueError(f"Error parsing file content: {str(e)}")
